#!/usr/bin/env python3
# SPDX-License-Identifier: GPL-3.0-or-later
# Copyright (C) 2026 James Sparenberg
"""
DocuBrowse PowerPoint extractor.

Extracts text and metadata from .pptx files using python-pptx.

Text extraction strategy:
  - Slide titles collected separately (used as document title if no core prop)
  - All text frames: paragraphs extracted in slide order
  - Tables: cells joined with ' | '
  - Notes pages skipped (usually presenter notes, not content)

Dependency: pip install python-pptx
"""

from pathlib import Path

_TEXT_LIMIT = 5000
_GENERIC_TITLES = frozenset({'powerpoint presentation', 'presentation', 'untitled'})


def extract_pptx(file_path: str) -> dict:
    result = {
        'title':       None,
        'author':      None,
        'subject':     None,
        'description': '',
        'text':        '',
        'snippet':     '',
        'doc_type':    'pptx',
        'success':     False,
        'error':       None,
    }
    try:
        from pptx import Presentation

        prs = Presentation(file_path)

        # Core properties
        cp = prs.core_properties
        raw_title = (getattr(cp, 'title', '') or '').strip()
        result['title'] = raw_title if raw_title.lower() not in _GENERIC_TITLES else None
        result['author']  = (getattr(cp, 'author',  '') or '').strip() or None
        result['subject'] = (getattr(cp, 'subject', '') or '').strip() or None
        comments = (getattr(cp, 'comments', '') or '').strip()
        keywords = (getattr(cp, 'keywords', '') or '').strip()
        result['description'] = comments or keywords or ''

        parts = []
        slide_titles = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    # Check for tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [
                                cell.text.strip()
                                for cell in row.cells
                                if cell.text.strip()
                            ]
                            if cells:
                                parts.append(' | '.join(cells))
                    continue

                # Collect slide title placeholder text separately
                try:
                    ph = shape.placeholder_format
                    is_title = ph is not None and ph.idx in (0, 1)
                except Exception:
                    is_title = False

                for para in shape.text_frame.paragraphs:  # noqa: E501
                    t = para.text.strip()
                    if not t:
                        continue
                    parts.append(t)
                    if is_title:
                        slide_titles.append(t)

        # Fall back to first slide title as document title
        if not result['title'] and slide_titles:
            result['title'] = slide_titles[0]

        full_text         = '\n'.join(parts)
        result['text']    = full_text[:_TEXT_LIMIT]
        result['snippet'] = full_text[:500]
        result['success'] = True
        return result

    except Exception as exc:
        result['error'] = str(exc)
        return result
